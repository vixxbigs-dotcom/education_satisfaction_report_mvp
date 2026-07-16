from __future__ import annotations

import io
import math
from pathlib import Path
from typing import Dict, List, Sequence, Tuple

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .models import ChoiceQuestion, ObjectiveQuestion, ReportData
from .slide_plan import SlideItem, build_slide_plan
from .text_utils import strip_leading_question_numbers

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ORANGE = RGBColor(237, 112, 44)
ORANGE_DARK = RGBColor(216, 92, 30)
DARK = RGBColor(68, 68, 68)
MUTED = RGBColor(111, 111, 111)
LIGHT = RGBColor(247, 247, 247)
LINE = RGBColor(215, 215, 215)
GRID = RGBColor(229, 229, 229)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 20, 20)
PURPLE = RGBColor(160, 43, 152)
BLUE = RGBColor(27, 159, 208)
GREEN = RGBColor(88, 169, 140)
GOLD = RGBColor(201, 161, 59)
FONT = "맑은 고딕"
PALETTE = [ORANGE, PURPLE, BLUE, RGBColor(243, 154, 102), RGBColor(124, 141, 181), GREEN, GOLD, RGBColor(155, 155, 155)]
ASSET_DIR = Path(__file__).resolve().parent.parent / "assets"
SECTION_BACKGROUND = ASSET_DIR / "multicampus_section_background.png"
LOGO_PATH = ASSET_DIR / "multicampus_logo.png"
SUBJECTIVE_PER_SLIDE = 5


def _set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def _add_text(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT,
              valign=MSO_ANCHOR.MIDDLE, margin=0.03, font_name=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin); tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text or "")
    run.font.name = font_name; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def _add_full_background(slide) -> None:
    if SECTION_BACKGROUND.exists():
        slide.shapes.add_picture(str(SECTION_BACKGROUND), 0, 0, SLIDE_W, SLIDE_H)
    else:
        _set_bg(slide, RGBColor(242, 242, 242))
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.25), 0, Inches(7.083), SLIDE_H)
        panel.fill.solid(); panel.fill.fore_color.rgb = WHITE; panel.line.fill.background()


def _add_logo(slide) -> None:
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(11.45), Inches(6.98), Inches(1.55), Inches(0.247))
    else:
        _add_text(slide, "multicampus", Inches(11.25), Inches(6.95), Inches(1.6), Inches(0.3), 10, True, BLACK, PP_ALIGN.RIGHT)


def _add_content_header(slide, main_title: str, section_no: str, section_title: str) -> None:
    _add_text(slide, main_title, Inches(0.32), Inches(0.27), Inches(6.4), Inches(0.56), 29, True, ORANGE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.6), Inches(0.73), Inches(7.75), Inches(0.035))
    line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.4), Inches(0.31), Inches(0.54), Inches(0.38))
    box.fill.solid(); box.fill.fore_color.rgb = ORANGE; box.line.fill.background()
    _add_text(slide, section_no, Inches(11.41), Inches(0.32), Inches(0.52), Inches(0.35), 16.5, False, WHITE, PP_ALIGN.CENTER)
    _add_text(slide, section_title, Inches(11.98), Inches(0.30), Inches(1.34), Inches(0.40), 15.5, True, ORANGE)
    _add_logo(slide)


def _style_table(table, header=True, first_col_gray=False, first_col_orange=False, body_font_size=11.0, header_font_size=11.5):
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08); cell.margin_top = Inches(0.035); cell.margin_bottom = Inches(0.035)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_header = r_idx == 0 and header
            is_orange = first_col_orange and c_idx == 0 and not is_header
            is_gray = first_col_gray and c_idx == 0 and not is_header
            cell.fill.solid(); cell.fill.fore_color.rgb = ORANGE if (is_header or is_orange) else LIGHT if is_gray else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if is_orange else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT; run.font.size = Pt(header_font_size if is_header else body_font_size)
                    run.font.color.rgb = WHITE if (is_header or is_orange) else DARK
                    run.font.bold = is_header or is_orange or is_gray


def _apply_table_row_heights(table, heights: Sequence[float]) -> None:
    for idx, height in enumerate(heights):
        if idx < len(table.rows):
            table.rows[idx].height = Inches(height)


def _add_cover(prs, report):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _add_full_background(slide)
    if report.company_name:
        _add_text(slide, report.company_name, Inches(6.9), Inches(2.18), Inches(5.25), Inches(0.36), 12, False, MUTED)
    _add_text(slide, report.course_name or "과정 이름", Inches(6.88), Inches(3.08), Inches(5.4), Inches(0.55), 29, True, ORANGE, valign=MSO_ANCHOR.BOTTOM)
    _add_text(slide, "결과보고서", Inches(6.88), Inches(3.62), Inches(5.4), Inches(0.58), 31, True, ORANGE, valign=MSO_ANCHOR.TOP)
    if report.schedule:
        _add_text(slide, report.schedule, Inches(6.9), Inches(4.55), Inches(4.8), Inches(0.32), 11, False, MUTED)


def _add_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _add_full_background(slide)
    _add_text(slide, "목차", Inches(6.88), Inches(2.48), Inches(4.8), Inches(0.55), 31, True, ORANGE)
    entries = [("01", "교육 개요", "교육 개요 · 커리큘럼"), ("02", "만족도 통계", "요약 · 공통문항 · 강사 · 선택형 · 주관식 · 시사점"), ("03", "현장 사진", "교육 현장 스케치")]
    for idx, (no, title, detail) in enumerate(entries):
        y = 3.32 + idx * 0.72
        _add_text(slide, no, Inches(6.9), Inches(y), Inches(0.55), Inches(0.28), 14, True, ORANGE)
        _add_text(slide, title, Inches(7.45), Inches(y), Inches(1.65), Inches(0.28), 13, True, DARK)
        _add_text(slide, detail, Inches(9.05), Inches(y), Inches(3.65), Inches(0.38), 8.5, False, MUTED, valign=MSO_ANCHOR.TOP)


def _add_section(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _add_full_background(slide)
    title = item.title.replace("1. ", "").replace("2. ", "").replace("3. ", "")
    _add_text(slide, item.payload.get("section_no", ""), Inches(6.9), Inches(2.7), Inches(1.2), Inches(0.38), 17, True, ORANGE)
    _add_text(slide, title, Inches(6.88), Inches(3.12), Inches(5.25), Inches(0.65), 31, True, ORANGE)
    _add_text(slide, " · ".join(item.payload.get("items", [])), Inches(6.9), Inches(4.03), Inches(5.25), Inches(0.55), 9.5, False, MUTED, valign=MSO_ANCHOR.TOP)


def _add_overview(prs, report):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, "1) 교육 개요", "01", "교육 개요")
    target = f"{report.target_text} {report.total_participants}명".strip() if report.total_participants else report.target_text
    rows = [("과정명", report.course_name), ("교육일정", report.schedule), ("교육방식", report.delivery_method), ("교육 대상", target), ("교육 목표", report.objective)]
    shape = slide.shapes.add_table(5, 2, Inches(0.46), Inches(1.32), Inches(12.08), Inches(5.55)); table = shape.table
    table.columns[0].width = Inches(2.2); table.columns[1].width = Inches(9.88)
    for i, (label, value) in enumerate(rows):
        table.cell(i, 0).text = label; table.cell(i, 1).text = str(value or "")
    _apply_table_row_heights(table, [1.02, 1.02, 1.02, 1.02, 1.35]); _style_table(table, header=False, first_col_orange=True, body_font_size=13)


def _add_curriculum(prs, report):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, "2) 커리큘럼", "01", "교육 개요")
    rows = report.curriculum[:8]; row_count = max(2, len(rows) + 1); heights = [0.48] + [0.5] * (row_count - 1)
    shape = slide.shapes.add_table(row_count, 4, Inches(0.62), Inches(1.45), Inches(12.08), Inches(sum(heights))); table = shape.table
    _apply_table_row_heights(table, heights)
    for idx, width in enumerate([1.15, 2.1, 6.75, 2.08]): table.columns[idx].width = Inches(width)
    for c, header in enumerate(["Day", "시간", "교육 내용", "강사/비고"]): table.cell(0, c).text = header
    if rows:
        for r, row in enumerate(rows, 1):
            for c, value in enumerate([row.day, row.time, row.content, row.instructor]): table.cell(r, c).text = str(value or "")
    else:
        table.cell(1, 0).merge(table.cell(1, 3)); table.cell(1, 0).text = ""
    _style_table(table, body_font_size=11.4, header_font_size=12)


def _all_questions(report):
    return list(report.objective_questions) + list(report.choice_questions) + list(report.subjective_questions)


def _add_survey_structure(prs, report, item):
    questions = _all_questions(report); start = int(item.payload.get("start", 0)); end = int(item.payload.get("end", len(questions))); chunk = questions[start:end]
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    title = "1) 설문 구성" + (f" ({int(item.payload.get('page_index', 0)) + 1})" if len(questions) > 8 else "")
    _add_content_header(slide, title, "02", "만족도 통계")
    row_count = max(2, len(chunk) + 1); heights = [0.48] + [0.55] * (row_count - 1)
    shape = slide.shapes.add_table(row_count, 2, Inches(0.62), Inches(1.45), Inches(12.08), Inches(sum(heights))); table = shape.table
    _apply_table_row_heights(table, heights); table.columns[0].width = Inches(2.45); table.columns[1].width = Inches(9.63)
    table.cell(0, 0).text = "항목"; table.cell(0, 1).text = "문항"
    for r, q in enumerate(chunk, 1):
        table.cell(r, 0).text = q.section_label; table.cell(r, 1).text = f"{start + r}. {strip_leading_question_numbers(q.question)}"
    if not chunk:
        table.cell(1, 0).merge(table.cell(1, 1)); table.cell(1, 0).text = ""
    _style_table(table, first_col_gray=True, body_font_size=11.2, header_font_size=12)


def _summary_label(q):
    if q.instructor_name: return f"{q.instructor_name}\n{q.instructor_metric or '만족도'}"
    return q.section_label


def _count_axis(max_count):
    raw = max(1, int(max_count))
    if raw <= 4:
        axis_max = 4
    elif raw <= 8:
        axis_max = 8
    elif raw <= 10:
        axis_max = 10
    else:
        axis_max = int(math.ceil(raw / 5.0) * 5)
    ticks = [axis_max * ratio / 4 for ratio in range(5)]
    return float(axis_max), ticks


def _format_count_tick(value):
    return str(int(value)) if float(value).is_integer() else f"{value:.1f}"


def _summary_suffix(item):
    page_count = int(item.payload.get("page_count", 1))
    if page_count <= 1:
        return ""
    page_index = int(item.payload.get("page_index", 0)) + 1
    return f" ({page_index}/{page_count})"


def _add_summary_distribution(prs, report, item, q):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    suffix = _summary_suffix(item)
    _add_content_header(slide, f"2) 만족도 요약{suffix}", "02", "만족도 통계")
    _add_text(slide, "■ 만족도 요약", Inches(.95), Inches(1.05), Inches(3.5), Inches(.38), 15, True, DARK)

    valid = max(0, int(q.valid_responses))
    if report.total_participants:
        count_line = (
            f"- 수강인원 총 {report.total_participants}명 중 {report.response_count}명 응답"
            f" ｜ 해당 문항 유효응답 {valid}명"
        )
    elif report.response_count != valid:
        count_line = f"- 총 {report.response_count}명 응답 ｜ 해당 문항 유효응답 {valid}명"
    else:
        count_line = f"- 유효 응답 {valid}명"
    _add_text(slide, count_line, Inches(1.18), Inches(1.39), Inches(8.8), Inches(.34), 13.2, False, BLACK)

    question_text = f"■ {q.section_label}  {strip_leading_question_numbers(q.question)}"
    _add_text(slide, question_text, Inches(.92), Inches(1.78), Inches(9.55), Inches(.54), _question_font_size(question_text), True, DARK, valign=MSO_ANCHOR.TOP)
    _add_text(slide, f"평균 {q.average:.1f}점", Inches(10.55), Inches(1.77), Inches(1.75), Inches(.42), 15, True, ORANGE, PP_ALIGN.RIGHT)

    max_count = max(max(q.counts) if q.counts else 0, 1)
    axis_max, ticks = _count_axis(max_count)
    chart_x, chart_y, chart_w, chart_h = 1.25, 2.62, 11.0, 3.45

    for tick in ticks:
        y = chart_y + chart_h - (tick / axis_max) * chart_h
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(chart_x), Inches(y), Inches(chart_w), Inches(.012))
        line.fill.solid(); line.fill.fore_color.rgb = GRID; line.line.fill.background()
        _add_text(slide, _format_count_tick(tick), Inches(.72), Inches(y - .13), Inches(.42), Inches(.25), 9, False, MUTED, PP_ALIGN.RIGHT)

    y_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(chart_x), Inches(chart_y), Inches(.012), Inches(chart_h))
    y_axis.fill.solid(); y_axis.fill.fore_color.rgb = RGBColor(135, 135, 135); y_axis.line.fill.background()
    x_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(chart_x), Inches(chart_y + chart_h), Inches(chart_w), Inches(.012))
    x_axis.fill.solid(); x_axis.fill.fore_color.rgb = RGBColor(135, 135, 135); x_axis.line.fill.background()

    values_counts = list(reversed(list(zip(q.scale_values, q.counts))))
    slot = chart_w / max(1, len(values_counts))
    bar_w = min(.52, slot * .58)
    for idx, (value, count) in enumerate(values_counts):
        center_x = chart_x + slot * idx + slot / 2
        bar_h = max(0.0, min(chart_h, count / axis_max * chart_h))
        bar_y = chart_y + chart_h - bar_h
        if count > 0:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(center_x - bar_w / 2),
                Inches(bar_y),
                Inches(bar_w),
                Inches(bar_h),
            )
            bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
        percentage = count / valid * 100 if valid else 0.0
        label_y = max(chart_y - .48, bar_y - .48) if count > 0 else chart_y + chart_h - .48
        _add_text(
            slide,
            f"{count}명\n({percentage:.0f}%)",
            Inches(center_x - slot * .48),
            Inches(label_y),
            Inches(slot * .96),
            Inches(.45),
            8.2,
            False,
            DARK,
            PP_ALIGN.CENTER,
            MSO_ANCHOR.BOTTOM,
            margin=0,
        )
        _add_text(slide, str(value), Inches(center_x - slot * .45), Inches(chart_y + chart_h + .06), Inches(slot * .9), Inches(.26), 10, False, DARK, PP_ALIGN.CENTER)

    _add_text(slide, "점수", Inches(6.15), Inches(6.48), Inches(.6), Inches(.24), 9, False, MUTED, PP_ALIGN.CENTER)


def _add_summary(prs, report, item):
    questions = [report.objective_questions[i] for i in item.payload.get("question_indices", [])]
    if item.payload.get("summary_mode") == "distribution" and questions:
        _add_summary_distribution(prs, report, item, questions[0])
        return

    scale_min = questions[0].scale_min if questions else 0
    scale_max = questions[0].scale_max if questions else 5
    span = max(1, scale_max - scale_min)
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    suffix = _summary_suffix(item)
    _add_content_header(slide, f"2) 만족도 요약{suffix}", "02", "만족도 통계")
    _add_text(slide, "■ 만족도 요약", Inches(0.95), Inches(1.05), Inches(3.5), Inches(0.38), 15, True, DARK)
    count_line = f"- 수강인원 총 {report.total_participants}명 중 {report.response_count}명 응답 ｜ 응답률 {report.response_rate}%" if report.total_participants else f"- 총 {report.response_count}명 응답"
    _add_text(slide, count_line, Inches(1.18), Inches(1.39), Inches(8.8), Inches(0.34), 14, False, BLACK)
    chart_x, chart_y, chart_w, chart_h = 1.4, 1.82, 11.35, 4.3
    for value in range(scale_min, scale_max + 1):
        ratio = (value - scale_min) / span
        y = chart_y + chart_h - ratio * chart_h
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(chart_x), Inches(y), Inches(chart_w), Inches(0.012)); line.fill.solid(); line.fill.fore_color.rgb = GRID; line.line.fill.background()
        _add_text(slide, f"{value:.1f}", Inches(0.78), Inches(y - 0.14), Inches(0.5), Inches(0.28), 10.5, False, MUTED, PP_ALIGN.RIGHT)
    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(chart_x), Inches(chart_y), Inches(0.012), Inches(chart_h)); axis.fill.solid(); axis.fill.fore_color.rgb = BLUE; axis.line.fill.background()
    n = max(1, len(questions)); slot = chart_w / n; bar_w = min(0.46, slot * 0.34)
    for idx, q in enumerate(questions):
        center_x = chart_x + slot * idx + slot / 2; bar_h = max(0, min(chart_h, ((q.average - scale_min) / span) * chart_h)); bar_y = chart_y + chart_h - bar_h
        if bar_h > 0:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(center_x - bar_w / 2), Inches(bar_y), Inches(bar_w), Inches(bar_h)); bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
        _add_text(slide, f"{q.average:.1f}", Inches(center_x - .35), Inches(bar_y - .3), Inches(.7), Inches(.28), 12.5, False, BLACK, PP_ALIGN.CENTER)
        _add_text(slide, _summary_label(q), Inches(center_x - slot * .48), Inches(chart_y + chart_h + .08), Inches(slot * .96), Inches(.68), 9.8, False, MUTED, PP_ALIGN.CENTER, MSO_ANCHOR.TOP)


def _axis_max(max_count):
    if max_count <= 9: return max(1, max_count), 1
    step = max(1, math.ceil(max_count / 9)); return int(math.ceil(max_count / step) * step), step


def _question_font_size(text):
    length = len(str(text or "")); return 14.5 if length <= 58 else 13.2 if length <= 90 else 12.2 if length <= 125 else 11.2


def _add_objective_horizontal(slide, q):
    raw_max = max(q.valid_responses, max(q.counts) if q.counts else 0, 1); axis_max, step = _axis_max(raw_max)
    chart_x, chart_y, chart_w, chart_h = 1.4, 2.02, 11.25, 4.13
    for tick in range(0, axis_max + 1, step):
        x = chart_x + chart_w * tick / axis_max
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(chart_y), Inches(.012), Inches(chart_h)); line.fill.solid(); line.fill.fore_color.rgb = GRID; line.line.fill.background()
        _add_text(slide, str(tick), Inches(x - .18), Inches(chart_y + chart_h + .06), Inches(.36), Inches(.26), 9.5, False, MUTED, PP_ALIGN.CENTER)
    row_h = chart_h / max(1, len(q.scale_labels))
    for idx, (label, count) in enumerate(zip(q.scale_labels, q.counts)):
        row_y = chart_y + idx * row_h
        _add_text(slide, label, Inches(chart_x + .08), Inches(row_y), Inches(3.2), Inches(.32), 13.2, False, DARK, valign=MSO_ANCHOR.BOTTOM)
        bar_y = row_y + .31; bar_h = min(.34, row_h * .42); bar_w = chart_w * count / axis_max if axis_max else 0
        if bar_w > 0:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(chart_x), Inches(bar_y), Inches(bar_w), Inches(bar_h)); bar.fill.solid(); bar.fill.fore_color.rgb = PALETTE[idx % len(PALETTE)]; bar.line.fill.background()
            if bar_w >= .35: _add_text(slide, str(count), Inches(chart_x + max(0, bar_w - .42)), Inches(bar_y), Inches(.35), Inches(bar_h), 11, True, WHITE, PP_ALIGN.RIGHT)
            else: _add_text(slide, str(count), Inches(chart_x + bar_w + .04), Inches(bar_y), Inches(.35), Inches(bar_h), 10.5, True, DARK)


def _add_objective_histogram(slide, q):
    chart_data = ChartData(); chart_data.categories = [str(v) for v in reversed(q.scale_values)]; chart_data.add_series("응답 수", list(reversed(q.counts)))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.05), Inches(11.25), Inches(4.35), chart_data).chart
    chart.has_legend = False; chart.has_title = False
    chart.value_axis.minimum_scale = 0; chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.name = FONT; chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.name = FONT; chart.value_axis.tick_labels.font.size = Pt(9)
    series = chart.series[0]; series.format.fill.solid(); series.format.fill.fore_color.rgb = ORANGE; series.format.line.fill.background()
    series.has_data_labels = True; labels = series.data_labels; labels.position = XL_LABEL_POSITION.OUTSIDE_END; labels.font.name = FONT; labels.font.size = Pt(9); labels.font.color.rgb = DARK


def _add_objective(prs, report, item):
    q = report.objective_questions[int(item.payload["question_index"])]
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    title = "3) 강사별 상세" if item.payload.get("group") == "lecturer" else "3) 공통문항"
    _add_content_header(slide, title, "02", "만족도 통계")
    display_order = int(item.payload.get("display_order", 1)); text = f"■ {q.section_label} {display_order}. {strip_leading_question_numbers(q.question)}"
    _add_text(slide, text, Inches(.88), Inches(1.15), Inches(10.2), Inches(.78), _question_font_size(text), True, DARK, valign=MSO_ANCHOR.TOP)
    _add_text(slide, f"평균 {q.average:.1f} / {q.scale_max}점", Inches(10.65), Inches(1.18), Inches(1.75), Inches(.38), 12.5, True, ORANGE, PP_ALIGN.RIGHT)
    if len(q.scale_values) >= 8: _add_objective_histogram(slide, q)
    else: _add_objective_horizontal(slide, q)


def _lecturer_metrics(report, names):
    metrics = []
    for name in names:
        for q in report.lecturers.get(name, []):
            metric = q.instructor_metric or "기타"
            if metric not in metrics: metrics.append(metric)
    return metrics


def _add_lecturer_comparison(prs, report, item):
    names = [report.instructors[i] for i in item.payload.get("lecturer_indices", [])]; metrics = _lecturer_metrics(report, names)
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, "3) 강사 비교", "02", "만족도 통계")
    _add_text(slide, "■ 강사별 문항 평균 비교", Inches(.78), Inches(1.12), Inches(5), Inches(.35), 15, True, DARK)
    row_count = max(2, len(names) + 1); col_count = max(2, len(metrics) + 2)
    heights = [.52] + [.7] * (row_count - 1)
    shape = slide.shapes.add_table(row_count, col_count, Inches(.72), Inches(1.65), Inches(11.9), Inches(sum(heights))); table = shape.table
    _apply_table_row_heights(table, heights); table.columns[0].width = Inches(2.0)
    metric_width = 8.0 / max(1, len(metrics))
    for idx in range(len(metrics)): table.columns[idx + 1].width = Inches(metric_width)
    table.columns[col_count - 1].width = Inches(1.9)
    table.cell(0, 0).text = "강사"
    for idx, metric in enumerate(metrics, 1): table.cell(0, idx).text = metric
    table.cell(0, col_count - 1).text = "종합 평균"
    for row_index, name in enumerate(names, 1):
        table.cell(row_index, 0).text = name
        metric_values: Dict[str, List[float]] = {}
        for question in report.lecturers.get(name, []):
            metric_values.setdefault(question.instructor_metric or "기타", []).append(question.average)
        by_metric = {metric: sum(values) / len(values) for metric, values in metric_values.items()}
        for col_index, metric in enumerate(metrics, 1): table.cell(row_index, col_index).text = f"{by_metric[metric]:.1f}" if metric in by_metric else "-"
        lecturer_questions = report.lecturers.get(name, [])
        total = sum(q.average for q in lecturer_questions) / len(lecturer_questions) if lecturer_questions else 0
        table.cell(row_index, col_count - 1).text = f"{total:.1f}"
    _style_table(table, first_col_gray=True, body_font_size=12, header_font_size=11.5)
    for row in range(1, row_count):
        cell = table.cell(row, col_count - 1); cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 247, 242)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs: run.font.bold = True; run.font.color.rgb = ORANGE


def _add_single_choice(prs, q: ChoiceQuestion, order: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, "4) 단일선택 문항", "02", "만족도 통계")
    text = f"■ {order}. {strip_leading_question_numbers(q.question)}"
    _add_text(slide, text, Inches(.75), Inches(1.13), Inches(11.6), Inches(.65), _question_font_size(text), True, DARK, valign=MSO_ANCHOR.TOP)
    data = ChartData(); data.categories = q.options; data.add_series("응답", q.counts)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(.85), Inches(1.9), Inches(5.3), Inches(4.65), data).chart
    chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.RIGHT; chart.legend.include_in_layout = False
    chart.legend.font.name = FONT; chart.legend.font.size = Pt(10)
    series = chart.series[0]; series.has_data_labels = True; labels = series.data_labels; labels.show_percentage = True; labels.show_category_name = False; labels.position = XL_LABEL_POSITION.BEST_FIT; labels.font.name = FONT; labels.font.size = Pt(9)
    for idx, point in enumerate(series.points):
        point.format.fill.solid(); point.format.fill.fore_color.rgb = PALETTE[idx % len(PALETTE)]
    y = 2.0
    for idx, (option, count, pct) in enumerate(zip(q.options[:9], q.counts[:9], q.percentages[:9])):
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(y + idx * .46), Inches(.18), Inches(.18)); marker.fill.solid(); marker.fill.fore_color.rgb = PALETTE[idx % len(PALETTE)]; marker.line.fill.background()
        _add_text(slide, option, Inches(7.28), Inches(y - .04 + idx * .46), Inches(3.15), Inches(.3), 10.5, False, DARK)
        _add_text(slide, f"{count}명 · {pct:.1f}%", Inches(10.45), Inches(y - .04 + idx * .46), Inches(1.55), Inches(.3), 10.5, True, DARK, PP_ALIGN.RIGHT)


def _add_multiple_choice(prs, q: ChoiceQuestion, order: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, "4) 복수선택 문항", "02", "만족도 통계")
    text = f"■ {order}. {strip_leading_question_numbers(q.question)}"
    _add_text(slide, text, Inches(.75), Inches(1.13), Inches(11.6), Inches(.65), _question_font_size(text), True, DARK, valign=MSO_ANCHOR.TOP)
    options = q.options[:10]; counts = q.counts[:10]; percentages = q.percentages[:10]
    max_count = max(counts) if counts else 1; top = 2.0; available = 4.45; row_h = available / max(1, len(options))
    for idx, (option, count, pct) in enumerate(zip(options, counts, percentages)):
        y = top + idx * row_h
        _add_text(slide, option, Inches(.72), Inches(y), Inches(3.0), Inches(row_h * .7), 10.5, False, DARK, PP_ALIGN.RIGHT)
        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(y + .07), Inches(6.5), Inches(min(.3, row_h * .45))); track.fill.solid(); track.fill.fore_color.rgb = GRID; track.line.fill.background()
        width = 6.5 * count / max_count if max_count else 0
        if width > 0:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(y + .07), Inches(width), Inches(min(.3, row_h * .45))); bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
        _add_text(slide, f"{count}명 · {pct:.1f}%", Inches(10.55), Inches(y), Inches(1.5), Inches(row_h * .7), 10.5, True, DARK)
    _add_text(slide, "※ 복수선택 문항은 비율 합계가 100%를 초과할 수 있습니다.", Inches(7.0), Inches(6.55), Inches(5.0), Inches(.25), 9, False, MUTED, PP_ALIGN.RIGHT)


def _add_choice(prs, report, item):
    index = int(item.payload["question_index"]); q = report.choice_questions[index]
    if q.selection_type == "multiple": _add_multiple_choice(prs, q, index + 1)
    else: _add_single_choice(prs, q, index + 1)


def _subjective_font_size(answers):
    total = sum(len(str(a)) for a in answers)
    return 13 if len(answers) <= 4 and total <= 420 else 12 if total <= 650 else 10.7


def _add_subjective(prs, report, item):
    index = int(item.payload["question_index"]); q = report.subjective_questions[index]; start = int(item.payload.get("chunk_index", 0)) * SUBJECTIVE_PER_SLIDE
    answers = q.answers[start:start + SUBJECTIVE_PER_SLIDE] or ["별도 의견 없음"]
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, f"5) 주관식 설문 결과{item.payload.get('page_suffix', '')}", "02", "만족도 통계")
    _add_text(slide, f"■ {q.section_label} {index + 1}. {strip_leading_question_numbers(q.question)}", Inches(.68), Inches(1.2), Inches(12), Inches(.68), 15.5, True, DARK, valign=MSO_ANCHOR.TOP)
    box = slide.shapes.add_textbox(Inches(.88), Inches(1.82), Inches(11.7), Inches(4.95)); tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(.02); tf.vertical_anchor = MSO_ANCHOR.TOP
    size = _subjective_font_size(answers)
    for idx, answer in enumerate(answers):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(6); p.line_spacing = 1.18
        run = p.add_run(); run.text = f"-  {answer}"; run.font.name = FONT; run.font.size = Pt(size); run.font.color.rgb = DARK


def _add_insights(prs, report):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, "6) 종합 시사점", "02", "만족도 통계")
    insights = report.insights or ["자동 분석된 시사점이 없습니다. 웹 입력부에서 내용을 추가해 주세요."]
    top = 1.55; box_h = min(1.05, 4.8 / max(1, len(insights)))
    for idx, text in enumerate(insights[:4]):
        y = top + idx * (box_h + .22)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.95), Inches(y), Inches(11.35), Inches(box_h)); box.fill.solid(); box.fill.fore_color.rgb = RGBColor(250, 250, 250); box.line.color.rgb = LINE
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.95), Inches(y), Inches(.12), Inches(box_h)); accent.fill.solid(); accent.fill.fore_color.rgb = ORANGE; accent.line.fill.background()
        _add_text(slide, str(idx + 1), Inches(1.2), Inches(y + .1), Inches(.55), Inches(box_h - .2), 20, True, ORANGE, PP_ALIGN.CENTER)
        _add_text(slide, text, Inches(1.85), Inches(y + .08), Inches(10.05), Inches(box_h - .16), 12.5, False, DARK)


def _fit_image(data, box_w, box_h):
    with Image.open(io.BytesIO(data)) as image: w, h = image.size
    ratio = min(box_w / w, box_h / h); return w * ratio, h * ratio


def _add_photos(prs, photos, item):
    page = int(item.payload.get("page_index", 0)); slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide); _add_content_header(slide, f"1) 현장 사진{item.payload.get('page_suffix', '')}", "03", "현장 사진")
    page_photos = photos[page * 6:(page + 1) * 6]
    left, top, gap_x, gap_y = .62, 1.4, .14, .14; box_w = (12.08 - gap_x * 2) / 3; box_h = (5.25 - gap_y) / 2
    for idx, (name, data) in enumerate(page_photos[:6]):
        row, col = divmod(idx, 3); x = left + col * (box_w + gap_x); y = top + row * (box_h + gap_y)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h)); bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(241, 241, 241); bg.line.color.rgb = LINE
        try:
            pw, ph = _fit_image(data, box_w, box_h); slide.shapes.add_picture(io.BytesIO(data), Inches(x + (box_w - pw) / 2), Inches(y + (box_h - ph) / 2), Inches(pw), Inches(ph))
        except Exception: _add_text(slide, name, Inches(x + .1), Inches(y + .1), Inches(box_w - .2), Inches(box_h - .2), 10.5, False, MUTED, PP_ALIGN.CENTER)


def _add_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _add_full_background(slide); _add_text(slide, "THANK YOU", Inches(6.88), Inches(3.08), Inches(5.1), Inches(.65), 34, True, ORANGE); _add_text(slide, "감사합니다", Inches(6.9), Inches(3.78), Inches(4.8), Inches(.38), 16, True, ORANGE)


def generate_pptx(report: ReportData, photos: Sequence[Tuple[str, bytes]] | None = None) -> bytes:
    photos = photos or []; prs = Presentation(); prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    while len(prs.slides): prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    for item in build_slide_plan(report, photo_count=len(photos)):
        if item.kind == "cover": _add_cover(prs, report)
        elif item.kind == "toc": _add_toc(prs)
        elif item.kind == "section": _add_section(prs, item)
        elif item.kind == "overview": _add_overview(prs, report)
        elif item.kind == "curriculum": _add_curriculum(prs, report)
        elif item.kind == "survey_structure": _add_survey_structure(prs, report, item)
        elif item.kind == "summary": _add_summary(prs, report, item)
        elif item.kind == "objective": _add_objective(prs, report, item)
        elif item.kind == "lecturer_comparison": _add_lecturer_comparison(prs, report, item)
        elif item.kind == "choice": _add_choice(prs, report, item)
        elif item.kind == "subjective": _add_subjective(prs, report, item)
        elif item.kind == "insights": _add_insights(prs, report)
        elif item.kind == "photos": _add_photos(prs, photos, item)
        elif item.kind == "thanks": _add_thanks(prs)
    output = io.BytesIO(); prs.save(output); return output.getvalue()
